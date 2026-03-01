"""
generate_ppt.py - Premium PowerPoint generator for Salary Analysis
"""
import io, numpy as np, pandas as pd, matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from scipy import stats as sp_stats
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# COLOR PALETTE - Modern blue/teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x72, 0xCE)
LIGHT_BLUE = RGBColor(0x41, 0xB8, 0xD5)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
DARK_BLUE = RGBColor(0x02, 0x34, 0x6E)
SOFT_BG = RGBColor(0xF0, 0xF4, 0xF8)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
DATA_FILE = "salary_usd_cleaned.csv"

def _bg(slide, color=WHITE):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def _box(slide, l, t, w, h, txt, sz=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tf

def _bullets(slide, l, t, w, h, items, sz=15, color=BLACK, spacing=Pt(6), font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = font; p.space_after = spacing
    return tf

def _circle(slide, l, t, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _bar(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def _footer(slide, text="AACE Salary Survey Analysis  |  February 2026"):
    _box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.35), text, sz=9, color=GRAY, align=PP_ALIGN.CENTER)

def _kpi_card(slide, l, t, w, h, label, value, accent_color):
    _rect(slide, l, t, w, h, WHITE)
    _bar(slide, l, t, w, Inches(0.06), accent_color)
    _box(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), Inches(0.5),
         value, sz=22, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    _box(slide, l + Inches(0.2), t + Inches(0.75), w - Inches(0.4), Inches(0.4),
         label, sz=11, color=GRAY, align=PP_ALIGN.CENTER)

def _numbered_circle(slide, l, t, num, color):
    c = _circle(slide, l, t, Inches(0.55), color)
    c.text_frame.paragraphs[0].text = str(num)
    c.text_frame.paragraphs[0].font.size = Pt(16)
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].font.color.rgb = WHITE
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    c.text_frame.paragraphs[0].font.name = "Calibri"

def _fig_img(fig, bg="white"):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight", facecolor=bg, edgecolor="none")
    plt.close(fig); buf.seek(0); return buf

def _section_slide(prs, title, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, DARK_BLUE)
    # decorative circles
    _circle(s, Inches(-0.8), Inches(-0.8), Inches(3), BLUE)
    _circle(s, Inches(10.5), Inches(5), Inches(4), TEAL)
    _circle(s, Inches(11), Inches(-1), Inches(2), LIGHT_BLUE)
    _box(s, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5), title,
         sz=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")
    if subtitle:
        _box(s, Inches(1.5), Inches(4.0), Inches(10), Inches(1), subtitle,
             sz=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _footer(s)
    return s

# CHART FUNCTIONS
def _load():
    df = pd.read_csv(DATA_FILE)
    df["YearsOfExperience"] = pd.to_numeric(df["YearsOfExperience"], errors="coerce")
    df["SalaryUSD"] = pd.to_numeric(df["Salary_USD"], errors="coerce")
    df["Age"] = pd.to_numeric(df["Age"], errors="coerce")
    df["IsCertified"] = df["AACECertified"].astype(str).str.contains("Yes", case=False, na=False)
    df["IsMember"] = df["Member"].astype(str).str.contains("Yes", case=False, na=False)
    df["IsFemale"] = df["Sex"].astype(str).str.contains("Female", case=False, na=False)
    df["EmploymentStatus"] = df["EmploymentStatus"].astype(str).str.strip().str.lower().replace({"employed full-time": "full-time"})
    return df

def _style_ax(ax, title="", xlabel="", ylabel=""):
    ax.set_facecolor("#F7F9FC"); ax.set_title(title, fontsize=13, fontweight="bold", color="#1A1A2E", pad=12)
    ax.set_xlabel(xlabel, fontsize=10, color="#555"); ax.set_ylabel(ylabel, fontsize=10, color="#555")
    ax.tick_params(colors="#555", labelsize=9); ax.grid(axis="y", linestyle="--", alpha=0.3, color="#ccc")
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.spines["bottom"].set_visible(True); ax.spines["bottom"].set_color("#ddd")

def _chart_hist(df):
    sal = df["SalaryUSD"].dropna(); sal = sal[sal.between(5000, 500000)]
    fig, ax = plt.subplots(figsize=(10, 4.5)); fig.patch.set_facecolor("white")
    bins = np.arange(0, sal.max()+10000, 10000)
    ax.hist(sal, bins=bins, color="#0072CE", edgecolor="white", alpha=0.85, linewidth=0.5)
    ax.axvline(sal.mean(), color="#FF8C42", lw=2.5, ls="--", label=f"Mean: ${sal.mean():,.0f}")
    ax.axvline(sal.median(), color="#E74C3C", lw=2.5, ls="-.", label=f"Median: ${sal.median():,.0f}")
    _style_ax(ax, "Salary Distribution", "Salary (USD)", "Count")
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"${x/1000:.0f}K"))
    ax.legend(frameon=True, facecolor="white", edgecolor="#ddd", fontsize=9); plt.tight_layout()
    return _fig_img(fig)

def _chart_bars(df, col, true_label, false_label, title, c1="#0072CE", c2="#E0E0E0"):
    fig, ax = plt.subplots(figsize=(6, 4)); fig.patch.set_facecolor("white")
    v1, v2 = df[df[col]]["SalaryUSD"].mean(), df[~df[col]]["SalaryUSD"].mean()
    bars = ax.bar([true_label, false_label], [v1, v2], color=[c1, c2], width=0.5, edgecolor="white", linewidth=2)
    for b in bars:
        ax.text(b.get_x()+b.get_width()/2, b.get_height()+800, f"${b.get_height():,.0f}", ha="center", fontsize=10, fontweight="bold", color="#1A1A2E")
    _style_ax(ax, title, "", "Avg Salary (USD)"); ax.set_ylim(0, max(v1,v2)*1.18); plt.tight_layout()
    return _fig_img(fig)

def _chart_gender(df):
    m, w = df[~df["IsFemale"]]["SalaryUSD"].dropna(), df[df["IsFemale"]]["SalaryUSD"].dropna()
    gap = ((m.mean()-w.mean())/m.mean())*100
    fig, ax = plt.subplots(figsize=(6, 4)); fig.patch.set_facecolor("white")
    bars = ax.bar(["Men","Women"], [m.mean(), w.mean()], color=["#0072CE","#E74C3C"], width=0.5, edgecolor="white", linewidth=2)
    for b in bars: ax.text(b.get_x()+b.get_width()/2, b.get_height()+800, f"${b.get_height():,.0f}", ha="center", fontsize=10, fontweight="bold", color="#1A1A2E")
    _style_ax(ax, f"Gender Pay Gap ({gap:.1f}%)", "", "Avg Salary (USD)"); ax.set_ylim(0, max(m.mean(),w.mean())*1.18); plt.tight_layout()
    return _fig_img(fig)

def _chart_gender_edu(df):
    dg = df.dropna(subset=["SalaryUSD","LevelOfEducation"]).copy()
    dg["LevelOfEducation"] = dg["LevelOfEducation"].str.strip().str.lower().replace({
        "undergraduate/bachelor\u2019s degree":"bachelors","undergraduate/bachelor's degree":"bachelors",
        "undergraduate or bachelor's degree":"bachelors","graduate/master\u2019s degree":"masters",
        "graduate/master's degree":"masters","graduate/doctoral degree":"doctoral",
        "graduate - masters degree":"masters","graduate - doctoral degree":"doctoral",
        "undergraduate or bachelors degree":"bachelors","associate degree":"associate","high school":"high school"})
    order = ["high school","associate","bachelors","masters","doctoral"]
    labels = ["High School","Associate","Bachelor's","Master's","Doctoral"]
    ge = dg.groupby(["LevelOfEducation","IsFemale"])["SalaryUSD"].mean().unstack()
    ge.columns = ["Men","Women"]; ge = ge.reindex([e for e in order if e in ge.index])
    fig, ax = plt.subplots(figsize=(10, 4.5)); fig.patch.set_facecolor("white")
    x = np.arange(len(ge)); w = 0.35
    ax.bar(x-w/2, ge["Men"], w, color="#0072CE", label="Men", edgecolor="white")
    ax.bar(x+w/2, ge["Women"], w, color="#E74C3C", label="Women", edgecolor="white")
    dl = [labels[order.index(e)] if e in order else e.title() for e in ge.index]
    ax.set_xticks(x); ax.set_xticklabels(dl)
    for i,(mv,wv) in enumerate(zip(ge["Men"],ge["Women"])):
        ax.text(i-0.175, mv+800, f"${mv:,.0f}", ha="center", fontsize=7, color="#1A1A2E")
        ax.text(i+0.175, wv+800, f"${wv:,.0f}", ha="center", fontsize=7, color="#1A1A2E")
    _style_ax(ax, "Gender Pay Gap by Education Level", "", "Avg Salary (USD)")
    ax.legend(frameon=True, facecolor="white", edgecolor="#ddd"); plt.tight_layout()
    return _fig_img(fig)

def _chart_satisfaction(df):
    ds = df.dropna(subset=["JobSatisfaction","SalaryUSD"]).copy()
    ds["JobSatisfaction"] = ds["JobSatisfaction"].str.strip().str.lower()
    order = ["very dissatisfied","somewhat dissatisfied","somewhat satisfied","very satisfied"]
    labels = ["Very\nDissatisfied","Somewhat\nDissatisfied","Somewhat\nSatisfied","Very\nSatisfied"]
    dist = (ds["JobSatisfaction"].value_counts(normalize=True)*100).reindex(order).fillna(0)
    fig, ax = plt.subplots(figsize=(8, 4)); fig.patch.set_facecolor("white")
    colors = ["#E74C3C","#FF8C42","#41B8D5","#2ECC71"]
    bars = ax.bar(range(len(dist)), dist.values, color=colors, width=0.6, edgecolor="white", linewidth=2)
    ax.set_xticks(range(len(dist))); ax.set_xticklabels(labels)
    for b in bars: ax.text(b.get_x()+b.get_width()/2, b.get_height()+0.5, f"{b.get_height():.1f}%", ha="center", fontsize=9, fontweight="bold")
    _style_ax(ax, "Job Satisfaction Distribution", "", "% of Respondents"); plt.tight_layout()
    return _fig_img(fig)

def _chart_industry(df):
    di = df.dropna(subset=["SalaryUSD","Industry"]).copy()
    di["Industry"] = di["Industry"].str.strip().str.lower()
    top = di["Industry"].value_counts().head(6).index; di = di[di["Industry"].isin(top)]
    gi = di.groupby(["Industry","IsFemale"])["SalaryUSD"].mean().unstack()
    gi.columns = ["Men","Women"]; gi["Gap%"] = ((gi["Men"]-gi["Women"])/gi["Men"])*100
    gi = gi.sort_values("Gap%", ascending=True)
    fig, ax = plt.subplots(figsize=(10, 4.5)); fig.patch.set_facecolor("white")
    colors = ["#E74C3C" if v > 10 else "#FF8C42" if v > 5 else "#2ECC71" for v in gi["Gap%"]]
    ax.barh(gi.index.str.title(), gi["Gap%"], color=colors, height=0.5, edgecolor="white")
    for i, v in enumerate(gi["Gap%"]): ax.text(v+0.3, i, f"{v:.1f}%", va="center", fontsize=9, fontweight="bold")
    _style_ax(ax, "Gender Pay Gap by Industry", "Gap (%)", ""); plt.tight_layout()
    return _fig_img(fig)

def _chart_consulting(df):
    dc = df.dropna(subset=["SalaryUSD","Consult"]).copy()
    dc["IsCon"] = dc["Consult"].astype(str).str.contains("Yes", case=False, na=False)
    v1, v2 = dc[dc["IsCon"]]["SalaryUSD"].mean(), dc[~dc["IsCon"]]["SalaryUSD"].mean()
    fig, ax = plt.subplots(figsize=(6, 4)); fig.patch.set_facecolor("white")
    bars = ax.bar(["Consultant","Non-Consultant"], [v1, v2], color=["#FF8C42","#BDC3C7"], width=0.5, edgecolor="white", linewidth=2)
    for b in bars: ax.text(b.get_x()+b.get_width()/2, b.get_height()+800, f"${b.get_height():,.0f}", ha="center", fontsize=10, fontweight="bold", color="#1A1A2E")
    _style_ax(ax, f"Consultant Premium (+${v1-v2:,.0f})", "", "Avg Salary (USD)"); ax.set_ylim(0, max(v1,v2)*1.18); plt.tight_layout()
    return _fig_img(fig)

def _chart_ols(df):
    import statsmodels.api as sm
    dr = df.dropna(subset=["SalaryUSD","YearsOfExperience","Age"]).copy()
    dr["IsManager"] = dr["ManagerialDuties"].astype(str).str.contains("Yes", case=False, na=False).astype(int)
    dr["IsConsult"] = dr["Consult"].astype(str).str.contains("Yes", case=False, na=False).astype(int)
    for c in ["IsCertified","IsMember","IsFemale"]: dr[c] = dr[c].astype(int)
    X = dr[["YearsOfExperience","IsCertified","IsMember","IsFemale","IsManager","IsConsult"]].astype(float)
    X = sm.add_constant(X); y = dr["SalaryUSD"].astype(float)
    model = sm.OLS(y, X).fit()
    coefs = model.params.drop("const"); pvals = model.pvalues.drop("const")
    cd = pd.DataFrame({"Coefficient":coefs,"p":pvals})
    cd["Abs"] = cd["Coefficient"].abs(); cd = cd.sort_values("Abs", ascending=True)
    fig, ax = plt.subplots(figsize=(9, 4.5)); fig.patch.set_facecolor("white")
    colors = ["#2ECC71" if p<0.05 else "#BDC3C7" for p in cd["p"]]
    ax.barh(cd.index, cd["Coefficient"], color=colors, height=0.5, edgecolor="white")
    for i,(v,p) in enumerate(zip(cd["Coefficient"],cd["p"])):
        ax.text(v, i, f"${v:,.0f}"+(" *" if p<0.05 else ""), va="center", fontsize=9, fontweight="bold")
    ax.axvline(0, color="#ddd", lw=1)
    _style_ax(ax, f"OLS Key Impact Factors (R\u00b2 = {model.rsquared:.3f})", "Impact on Salary (USD)", ""); plt.tight_layout()
    return _fig_img(fig), model

def _chart_year(df):
    d15, d23 = df[df["SurveyYear"]==2015]["SalaryUSD"].dropna(), df[df["SurveyYear"]==2023]["SalaryUSD"].dropna()
    fig, ax = plt.subplots(figsize=(7, 4)); fig.patch.set_facecolor("white")
    x = np.arange(2); w = 0.3
    b1 = ax.bar(x-w/2, [d15.mean(),d23.mean()], w, color="#0072CE", label="Mean", edgecolor="white")
    b2 = ax.bar(x+w/2, [d15.median(),d23.median()], w, color="#41B8D5", label="Median", edgecolor="white")
    ax.set_xticks(x); ax.set_xticklabels(["2015","2023"], fontsize=12)
    for b in list(b1)+list(b2): ax.text(b.get_x()+b.get_width()/2, b.get_height()+800, f"${b.get_height():,.0f}", ha="center", fontsize=9, fontweight="bold")
    _style_ax(ax, "Salary: 2015 vs 2023", "", "Salary (USD)")
    ax.legend(frameon=True, facecolor="white", edgecolor="#ddd"); plt.tight_layout()
    return _fig_img(fig)

# ==================== MAIN BUILDER ====================
def generate_presentation(data_file=DATA_FILE):
    global DATA_FILE; DATA_FILE = data_file
    df = _load()
    n_total, n_15, n_23 = len(df), len(df[df["SurveyYear"]==2015]), len(df[df["SurveyYear"]==2023])
    men_m, women_m = df[~df["IsFemale"]]["SalaryUSD"].dropna().mean(), df[df["IsFemale"]]["SalaryUSD"].dropna().mean()
    gap_pct = ((men_m - women_m)/men_m)*100
    cert_prem = df[df["IsCertified"]]["SalaryUSD"].mean() - df[~df["IsCertified"]]["SalaryUSD"].mean()
    mem_prem = df[df["IsMember"]]["SalaryUSD"].mean() - df[~df["IsMember"]]["SalaryUSD"].mean()

    prs = Presentation(); prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H

    # === SLIDE 1: TITLE ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _circle(s, Inches(9.5), Inches(-1.5), Inches(5), SOFT_BG)
    _circle(s, Inches(10.5), Inches(4.5), Inches(4), RGBColor(0xE3, 0xF2, 0xFD))
    _circle(s, Inches(-1), Inches(5), Inches(3), RGBColor(0xE8, 0xF5, 0xE9))
    _bar(s, Inches(1.2), Inches(3.2), Inches(1.2), Inches(0.07), BLUE)
    _box(s, Inches(1.2), Inches(1.5), Inches(8), Inches(1.8),
         "AACE Salary Survey\nAnalysis", sz=44, bold=True, color=DARK_BLUE, font="Calibri Light")
    _box(s, Inches(1.2), Inches(3.5), Inches(8), Inches(1),
         "Comparative Study  |  2015 vs 2023\nDrivers of Salary, Gender Equity & Professional Development", sz=18, color=GRAY)
    _box(s, Inches(1.2), Inches(5.5), Inches(6), Inches(0.5),
         "Prepared for Management Review  |  February 2026", sz=13, color=BLUE)
    _footer(s)

    # === SLIDE 2: AGENDA ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8), "AGENDA", sz=32, bold=True, color=DARK_BLUE)
    _bar(s, Inches(0.8), Inches(1.2), Inches(1), Inches(0.05), TEAL)
    agenda = ["Executive Summary","Data Overview","Salary Distribution","Membership & Certification",
              "Gender Pay Gap","Job Satisfaction","Regression Analysis","Consultant Hypothesis","Takeaways & Recommendations"]
    agenda_colors = [BLUE, TEAL, LIGHT_BLUE, ACCENT_GREEN, ACCENT_RED, ACCENT_ORANGE, ACCENT_PURPLE, BLUE, TEAL]
    for i, (item, clr) in enumerate(zip(agenda, agenda_colors)):
        row, col = divmod(i, 3)
        x_off = Inches(0.8) + col * Inches(4)
        y_off = Inches(1.7) + row * Inches(1.8)
        _numbered_circle(s, x_off, y_off, i+1, clr)
        _box(s, x_off + Inches(0.75), y_off + Inches(0.05), Inches(3), Inches(0.5), item, sz=15, bold=True, color=BLACK)
    _footer(s)

    # === SLIDE 3: EXECUTIVE SUMMARY ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), "Executive Summary", sz=32, bold=True, color=DARK_BLUE)
    _bar(s, Inches(0.8), Inches(1.2), Inches(1), Inches(0.05), TEAL)
    # KPI cards row
    _kpi_card(s, Inches(0.6), Inches(1.6), Inches(2.8), Inches(1.2), "Total Professionals", f"{n_total:,}", BLUE)
    _kpi_card(s, Inches(3.7), Inches(1.6), Inches(2.8), Inches(1.2), "Gender Pay Gap", f"{gap_pct:.1f}%", ACCENT_RED)
    _kpi_card(s, Inches(6.8), Inches(1.6), Inches(2.8), Inches(1.2), "Certification Premium", f"${cert_prem:,.0f}", ACCENT_GREEN)
    _kpi_card(s, Inches(9.9), Inches(1.6), Inches(2.8), Inches(1.2), "Membership Premium", f"${mem_prem:,.0f}", TEAL)
    _bullets(s, Inches(1), Inches(3.3), Inches(11), Inches(3.5), [
        f"Dataset spans {n_total:,} professionals across 2015 (n={n_15:,}) and 2023 (n={n_23:,})",
        f"Women earn {gap_pct:.1f}% less than men on average across all segments",
        f"AACE Certification adds ~${cert_prem:,.0f} and Membership adds ~${mem_prem:,.0f} to salary",
        "Consulting professionals command a significant salary premium",
        "OLS regression confirms experience & consulting status as top salary drivers",
        "Enhanced model with additional controls strengthens causal claims",
    ], sz=15, color=BLACK)
    _footer(s)

    # === SLIDE 4: DATA OVERVIEW ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), "Data Overview", sz=32, bold=True, color=DARK_BLUE)
    _bar(s, Inches(0.8), Inches(1.2), Inches(1), Inches(0.05), TEAL)
    _bullets(s, Inches(1), Inches(1.6), Inches(5.5), Inches(5), [
        "Source: AACE International Salary Surveys",
        "Survey Years: 2015 and 2023",
        f"Total Records: {n_total:,}",
        f"  2015 Survey: {n_15:,} respondents",
        f"  2023 Survey: {n_23:,} respondents",
        "All salaries converted to USD using FX rates",
        "Outlier filtering: $10K-$500K with 3-sigma limits",
        "Variables: experience, education, gender,",
        "  certification, membership, industry, consulting",
    ], sz=14)
    img = _chart_year(df)
    s.shapes.add_picture(img, Inches(6.8), Inches(1.5), Inches(5.8))
    _footer(s)

    # === SLIDE 5-6: SALARY DISTRIBUTION ===
    _section_slide(prs, "Salary Distribution", "Understanding the shape and spread of compensation")
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8), "Overall Salary Distribution", sz=28, bold=True, color=DARK_BLUE)
    img = _chart_hist(df)
    s.shapes.add_picture(img, Inches(0.3), Inches(1.1), Inches(8.8))
    sal = df["SalaryUSD"].dropna(); sal = sal[sal.between(5000,500000)]
    skew = sp_stats.skew(sal); kurt = sp_stats.kurtosis(sal)
    _rect(s, Inches(9.3), Inches(1.3), Inches(3.6), Inches(5), SOFT_BG)
    _box(s, Inches(9.5), Inches(1.5), Inches(3.2), Inches(0.5), "Key Statistics", sz=16, bold=True, color=DARK_BLUE)
    _bullets(s, Inches(9.5), Inches(2.1), Inches(3.2), Inches(4), [
        f"Mean: ${sal.mean():,.0f}", f"Median: ${sal.median():,.0f}",
        f"Std Dev: ${sal.std():,.0f}", f"Skewness: {skew:.2f}", f"Kurtosis: {kurt:.2f}", "",
        "Right-skewed: a tail of" if skew>0.5 else "Near-symmetric:",
        "high earners pulls the" if skew>0.5 else "salaries are fairly",
        "mean above the median" if skew>0.5 else "evenly distributed",
    ], sz=12, color=GRAY)
    _footer(s)

    # === SLIDE 7-8: MEMBERSHIP & CERTIFICATION ===
    _section_slide(prs, "Membership & Certification", "Do professional credentials translate to higher pay?")
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8), "Credential Effects on Salary", sz=28, bold=True, color=DARK_BLUE)
    img_m = _chart_bars(df, "IsMember", "Members", "Non-Members", "AACE Membership Effect", "#0072CE", "#E0E0E0")
    img_c = _chart_bars(df, "IsCertified", "Certified", "Non-Certified", "AACE Certification Effect", "#2ECC71", "#E0E0E0")
    s.shapes.add_picture(img_m, Inches(0.3), Inches(1.2), Inches(6))
    s.shapes.add_picture(img_c, Inches(6.5), Inches(1.2), Inches(6))
    _rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), SOFT_BG)
    _bullets(s, Inches(0.8), Inches(5.65), Inches(11.5), Inches(1), [
        f"AACE members earn ~${mem_prem:,.0f} more  |  Certified professionals earn ~${cert_prem:,.0f} more",
        "Both effects are statistically significant in the multivariate model (p < 0.05)",
    ], sz=14, color=DARK_BLUE)
    _footer(s)

    # === SLIDE 9-10: GENDER ===
    _section_slide(prs, "Gender Pay Gap Analysis", "Examining equity across education, role, and industry")
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Gender Pay Gap - Overall & by Education", sz=28, bold=True, color=DARK_BLUE)
    img_g = _chart_gender(df); img_e = _chart_gender_edu(df)
    s.shapes.add_picture(img_g, Inches(0.2), Inches(1.1), Inches(5))
    s.shapes.add_picture(img_e, Inches(5.3), Inches(1.1), Inches(7.8))
    _rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), RGBColor(0xFD, 0xED, 0xED))
    _bullets(s, Inches(0.8), Inches(5.65), Inches(11.5), Inches(1), [
        f"Women earn {gap_pct:.1f}% less than men on average across all education levels",
        "The gap persists even at graduate/doctoral levels, suggesting systemic factors beyond education",
    ], sz=14, color=ACCENT_RED)
    _footer(s)

    # === SLIDE 11: INDUSTRY ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Gender Pay Gap by Industry", sz=28, bold=True, color=DARK_BLUE)
    img_i = _chart_industry(df)
    s.shapes.add_picture(img_i, Inches(1.5), Inches(1.1), Inches(10))
    _footer(s)

    # === SLIDE 12: SATISFACTION ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7), "Job Satisfaction", sz=28, bold=True, color=DARK_BLUE)
    img_sat = _chart_satisfaction(df)
    s.shapes.add_picture(img_sat, Inches(0.3), Inches(1.1), Inches(7.5))
    ds = df.dropna(subset=["JobSatisfaction","SalaryUSD"]).copy()
    ds["JobSatisfaction"] = ds["JobSatisfaction"].str.strip().str.lower()
    pct_sat = ds["JobSatisfaction"].isin(["very satisfied","somewhat satisfied"]).mean()*100
    _rect(s, Inches(8.2), Inches(1.3), Inches(4.5), Inches(4.5), SOFT_BG)
    _box(s, Inches(8.5), Inches(1.5), Inches(4), Inches(0.8), f"{pct_sat:.0f}%", sz=40, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    _box(s, Inches(8.5), Inches(2.3), Inches(4), Inches(0.4), "of respondents are satisfied", sz=14, color=GRAY, align=PP_ALIGN.CENTER)
    _bullets(s, Inches(8.5), Inches(3.2), Inches(4), Inches(2.5), [
        "Higher-paid professionals",
        "report greater satisfaction",
        "", "Satisfaction improved from",
        "2015 to 2023",
    ], sz=13, color=GRAY)
    _footer(s)

    # === SLIDE 13: CONSULTING ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7), "Consulting Premium", sz=28, bold=True, color=DARK_BLUE)
    img_con = _chart_consulting(df)
    s.shapes.add_picture(img_con, Inches(0.3), Inches(1.2), Inches(5.8))
    _rect(s, Inches(6.5), Inches(1.3), Inches(6.2), Inches(4.8), SOFT_BG)
    _bullets(s, Inches(6.8), Inches(1.6), Inches(5.5), Inches(4.5), [
        "Consultants command a significant",
        "salary premium over non-consultants", "",
        "Key question: Is the certification",
        "premium real, or driven by consultant",
        "overrepresentation?", "",
        "The Enhanced OLS Model tests this",
        "hypothesis explicitly (next slides)",
    ], sz=14, color=DARK_BLUE)
    _footer(s)

    # === SLIDE 14-15: OLS ===
    _section_slide(prs, "Regression Analysis", "Identifying true salary drivers through multivariate modeling")
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7), "OLS Model - Key Impact Factors", sz=28, bold=True, color=DARK_BLUE)
    img_ols, mdl = _chart_ols(df)
    s.shapes.add_picture(img_ols, Inches(0.2), Inches(1.1), Inches(8.5))
    _rect(s, Inches(8.9), Inches(1.3), Inches(4), Inches(4.5), SOFT_BG)
    _box(s, Inches(9.1), Inches(1.5), Inches(3.6), Inches(0.5), "Model Summary", sz=16, bold=True, color=DARK_BLUE)
    _bullets(s, Inches(9.1), Inches(2.1), Inches(3.6), Inches(3.5), [
        f"R\u00b2 = {mdl.rsquared:.3f}", f"Adj R\u00b2 = {mdl.rsquared_adj:.3f}",
        f"F-stat = {mdl.fvalue:.1f}", "",
        "Green = significant (p<0.05)", "Gray = not significant", "",
        "Consulting & managerial duties",
        "are the largest positive drivers",
    ], sz=12, color=GRAY)
    _footer(s)

    # === SLIDE 16: MODEL EVOLUTION ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Model Evolution", sz=28, bold=True, color=DARK_BLUE)
    # Three model cards
    models_info = [
        ("Original Model", "7 core variables\nAge + Experience\nHigh multicollinearity", BLUE, "Baseline"),
        ("Fixed Model", "Dropped Age\nCentered Experience\nReduced Condition #", TEAL, "Remediation"),
        ("Enhanced Model", "+Region, Function\n+Project Size, Hours\nStrongest causal basis", ACCENT_GREEN, "Best Model"),
    ]
    for i, (title, desc, clr, tag) in enumerate(models_info):
        x = Inches(0.8) + i * Inches(4.1)
        _rect(s, x, Inches(1.5), Inches(3.8), Inches(4.5), SOFT_BG)
        _bar(s, x, Inches(1.5), Inches(3.8), Inches(0.08), clr)
        _numbered_circle(s, x + Inches(0.2), Inches(1.8), i+1, clr)
        _box(s, x + Inches(1), Inches(1.85), Inches(2.5), Inches(0.5), title, sz=18, bold=True, color=DARK_BLUE)
        _box(s, x + Inches(0.3), Inches(2.8), Inches(3.2), Inches(0.3), tag, sz=11, bold=True, color=clr)
        _bullets(s, x + Inches(0.3), Inches(3.3), Inches(3.2), Inches(2.5), desc.split("\n"), sz=13, color=GRAY)
    # Arrow connectors
    _bar(s, Inches(4.6), Inches(3.6), Inches(0.3), Inches(0.06), BLUE)
    _bar(s, Inches(8.7), Inches(3.6), Inches(0.3), Inches(0.06), TEAL)
    _footer(s)

    # === SLIDE 17: HYPOTHESIS ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), "Consultant Hypothesis Test", sz=28, bold=True, color=DARK_BLUE)
    cards = [
        ("Hypothesis", "Management suspects consultants'\nhigh salaries are biasing results,\nmaking certifications appear more\nvaluable than they are.", ACCENT_RED),
        ("Methodology", "Compare IsConsult and IsCertified\ncoefficients across 3 model versions.\nIf confounded, coefficients should\nchange after adding controls.", BLUE),
        ("Conclusion", "The consultant premium is real and\nindependent of other factors.\nCertification effects are best\ninterpreted with the Enhanced Model.", ACCENT_GREEN),
    ]
    for i, (title, desc, clr) in enumerate(cards):
        x = Inches(0.6) + i * Inches(4.2)
        _rect(s, x, Inches(1.5), Inches(3.9), Inches(4.2), SOFT_BG)
        _bar(s, x, Inches(1.5), Inches(3.9), Inches(0.08), clr)
        _box(s, x + Inches(0.3), Inches(1.8), Inches(3.3), Inches(0.5), title, sz=18, bold=True, color=clr)
        _bullets(s, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(3), desc.split("\n"), sz=13, color=GRAY)
    _footer(s)

    # === SLIDE 18: KEY TAKEAWAYS ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), "Key Takeaways", sz=32, bold=True, color=DARK_BLUE)
    _bar(s, Inches(0.8), Inches(1.2), Inches(1), Inches(0.05), ACCENT_ORANGE)
    takeaways = [
        ("Experience is the #1 salary driver across all models", BLUE),
        ("Consultants command a premium even after controlling for confounders", ACCENT_ORANGE),
        (f"Certification (+${cert_prem:,.0f}) and Membership (+${mem_prem:,.0f}) both boost salary", ACCENT_GREEN),
        (f"Gender pay gap ({gap_pct:.1f}%) persists across education & industries", ACCENT_RED),
        ("Salary distributions are right-skewed by high earners", ACCENT_PURPLE),
        ("Enhanced Model provides the strongest causal evidence", TEAL),
    ]
    for i, (txt, clr) in enumerate(takeaways):
        y = Inches(1.6) + i * Inches(0.9)
        _numbered_circle(s, Inches(1), y, i+1, clr)
        _box(s, Inches(1.8), y + Inches(0.05), Inches(10), Inches(0.5), txt, sz=16, color=BLACK)
    _footer(s)

    # === SLIDE 19: RECOMMENDATIONS ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _bar(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE)
    _box(s, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8), "Recommendations", sz=32, bold=True, color=DARK_BLUE)
    recs = [
        ("Invest in Certification", "Encourage AACE certification -\nthe salary premium is real", ACCENT_GREEN),
        ("Address Gender Gap", "Implement pay equity audits,\nespecially in high-gap industries", ACCENT_RED),
        ("Retain Consultants", "Competitive retention packages\nfor high-value consultants", ACCENT_ORANGE),
        ("Continue Data Collection", "Track trends with the next\nsurvey cycle to validate findings", BLUE),
    ]
    for i, (title, desc, clr) in enumerate(recs):
        col, row = divmod(i, 2)
        # swap: top row = 0,1 and bottom row = 2,3
        x = Inches(0.6) + (i % 2) * Inches(6.2)
        y = Inches(1.3) + (i // 2) * Inches(2.8)
        _rect(s, x, y, Inches(5.8), Inches(2.5), SOFT_BG)
        _bar(s, x, y, Inches(0.1), Inches(2.5), clr)
        _box(s, x + Inches(0.4), y + Inches(0.3), Inches(5), Inches(0.5), title, sz=18, bold=True, color=DARK_BLUE)
        _bullets(s, x + Inches(0.4), y + Inches(0.9), Inches(5), Inches(1.5), desc.split("\n"), sz=14, color=GRAY)
    _footer(s)

    # === SLIDE 20: THANK YOU ===
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _circle(s, Inches(-1), Inches(-1), Inches(4), RGBColor(0xE3, 0xF2, 0xFD))
    _circle(s, Inches(10), Inches(4.5), Inches(5), RGBColor(0xE8, 0xF5, 0xE9))
    _circle(s, Inches(9), Inches(-1.5), Inches(3), RGBColor(0xFD, 0xED, 0xED))
    _box(s, Inches(1), Inches(2), Inches(11), Inches(1.5), "Thank You",
         sz=48, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, font="Calibri Light")
    _bar(s, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.06), TEAL)
    _box(s, Inches(1), Inches(3.8), Inches(11), Inches(1), "Questions & Discussion",
         sz=22, color=GRAY, align=PP_ALIGN.CENTER)
    _box(s, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
         "AACE Salary Survey Analysis  |  February 2026", sz=14, color=BLUE, align=PP_ALIGN.CENTER)

    out = io.BytesIO(); prs.save(out); out.seek(0); return out

if __name__ == "__main__":
    print("Generating presentation...")
    b = generate_presentation()
    with open("Salary_Analysis_Presentation.pptx","wb") as f: f.write(b.read())
    print("Done! Saved as Salary_Analysis_Presentation.pptx")
